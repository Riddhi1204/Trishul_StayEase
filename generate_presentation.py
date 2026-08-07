import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6] # Blank slide

    # Palette
    DARK_BG = RGBColor(15, 30, 22)      # #0F1E16
    CARD_BG = RGBColor(24, 48, 36)      # #183024
    LIGHT_BG = RGBColor(248, 250, 248)  # #F8FAF8
    PRIMARY = RGBColor(45, 106, 79)     # #2D6A4F
    ACCENT = RGBColor(82, 183, 136)     # #52B788
    GOLD = RGBColor(234, 179, 8)        # #EAB308
    TEXT_DARK = RGBColor(30, 41, 59)    # #1E293B
    TEXT_MUTED = RGBColor(100, 116, 139)# #64748B
    WHITE = RGBColor(255, 255, 255)
    CARD_LIGHT = RGBColor(255, 255, 255)
    BORDER_COLOR = RGBColor(226, 232, 240)

    def set_slide_background(slide, color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_header(slide, section_num, title, subtitle=None):
        # Section pill
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.5), Inches(2.2), Inches(0.35))
        pill.fill.solid()
        pill.fill.fore_color.rgb = PRIMARY
        pill.line.color.rgb = ACCENT
        pill.line.width = Pt(1)
        p_tf = pill.text_frame
        p_tf.word_wrap = True
        p = p_tf.paragraphs[0]
        p.text = f"SECTION {section_num}"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Title
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.9), Inches(11.7), Inches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = DARK_BG

        if subtitle:
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.font.size = Pt(13)
            p2.font.color.rgb = TEXT_MUTED

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide (Dark Theme)
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, DARK_BG)

    # Accent decorative glow bar
    bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(0.15), Inches(5.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    # Main Title text box
    tb = slide1.shapes.add_textbox(Inches(1.2), Inches(1.2), Inches(11.0), Inches(2.2))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Trishul StayEase 🏡"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    p2 = tf.add_paragraph()
    p2.text = "Direct Booking Engine for Eco-Homestays & Sustainable Tourism"
    p2.font.size = Pt(18)
    p2.font.color.rgb = ACCENT

    p3 = tf.add_paragraph()
    p3.text = "Final Internship Evaluation Presentation • Full Stack Development (WD-05)"
    p3.font.size = Pt(13)
    p3.font.color.rgb = RGBColor(180, 205, 195)

    # Student Info Card
    card1 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(3.8), Inches(10.8), Inches(2.4))
    card1.fill.solid()
    card1.fill.fore_color.rgb = CARD_BG
    card1.line.color.rgb = PRIMARY
    card1.line.width = Pt(1.5)

    ctf = card1.text_frame
    ctf.word_wrap = True
    ctf.margin_left = Inches(0.4)
    ctf.margin_top = Inches(0.3)

    items = [
        ("Candidate Name:", "Riddhi Kumari"),
        ("Intern ID:", "26100462"),
        ("University / College:", "[Your University / College Name]"),
        ("Project Domain:", "Full Stack Web Development (FastAPI + React + MongoDB Atlas + Gemini AI)"),
        ("Live Demo:", "https://trishul-stay-ease.vercel.app")
    ]
    for i, (k, v) in enumerate(items):
        p = ctf.paragraphs[0] if i == 0 else ctf.add_paragraph()
        p.text = f"{k} "
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = ACCENT
        
        run = p.add_run()
        run.text = v
        run.font.bold = False
        run.font.color.rgb = WHITE

    # -------------------------------------------------------------
    # SLIDE 2: Problem Statement
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2, LIGHT_BG)
    add_header(slide2, "1 & 2", "What Problem Are We Solving?", "Addressing challenges faced by eco-homestay owners and sustainable travellers")

    problems = [
        ("💸 High OTA Commission Rates", [
            "Traditional travel platforms charge 15% - 25% commission fees per booking.",
            "Heavily reduces profit margins for small, family-run rural homestay owners.",
            "Forces eco-hosts to inflate prices for end guests."
        ]),
        ("🔍 Fragmented Eco-Discovery", [
            "Mindful travellers struggle to find authentic, verified sustainable stays.",
            "Greenwashing makes it difficult to verify authentic eco-practices.",
            "Lack of dedicated filters for local cultural & eco-activities."
        ]),
        ("⚡ Complex Direct Booking & Planning", [
            "Small homestays lack digital infrastructure for direct, commission-free reservations.",
            "Manual scheduling leads to double bookings and slow response times.",
            "Travellers lack personalized, budget-friendly sustainable itineraries."
        ])
    ]

    for idx, (title, bullets) in enumerate(problems):
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + idx*3.95), Inches(1.8), Inches(3.8), Inches(5.0))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_LIGHT
        card.line.color.rgb = BORDER_COLOR
        card.line.width = Pt(1)

        ctf = card.text_frame
        ctf.word_wrap = True
        ctf.margin_left = Inches(0.25)
        ctf.margin_right = Inches(0.25)
        ctf.margin_top = Inches(0.3)

        p = ctf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = PRIMARY
        p.space_after = Pt(12)

        for b in bullets:
            pb = ctf.add_paragraph()
            pb.text = f"•  {b}"
            pb.font.size = Pt(11.5)
            pb.font.color.rgb = TEXT_DARK
            pb.space_after = Pt(8)

    # -------------------------------------------------------------
    # SLIDE 3: Tech Stack Selected & Justifications
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3, LIGHT_BG)
    add_header(slide3, "3", "Tech Stack Selected Along with Reasons", "Modern, scalable, and type-safe architecture for maximum developer productivity")

    tech_stack = [
        ("🎨 Frontend: React 19 + Vite", "UI & Client Architecture", [
            "Lightning-fast Hot Module Replacement (HMR) & sub-second build times.",
            "Component-driven modularity with React Router 7 for fluid SPA navigation.",
            "Context API for state management (Auth, Dark/Light theme)."
        ]),
        ("⚡ Backend: FastAPI (Python 3.11)", "High-Performance ASGI REST API", [
            "Asynchronous execution handling high-concurrency requests seamlessly.",
            "Pydantic v2 for automatic request/response data validation and sanitization.",
            "Auto-generated interactive Swagger & OpenAPI documentation."
        ]),
        ("🗄️ Database: MongoDB Atlas + Motor", "Cloud Document Database", [
            "Motor async driver provides non-blocking I/O queries for the FastAPI event loop.",
            "Flexible JSON schema accommodates dynamic amenities, image arrays & reviews.",
            "Managed cloud high availability with zero server maintenance."
        ]),
        ("🤖 AI: Google Gemini 1.5 Flash", "Intelligent Itinerary Generator", [
            "State-of-the-art inference speed and structured JSON output guarantee.",
            "Built-in server-side TTL caching reduces redundant API overhead."
        ])
    ]

    for idx, (title, category, bullets) in enumerate(tech_stack):
        col = idx % 2
        row = idx // 2
        card = slide3.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, 
            Inches(0.8 + col*5.95), 
            Inches(1.8 + row*2.6), 
            Inches(5.75), 
            Inches(2.4)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_LIGHT
        card.line.color.rgb = BORDER_COLOR
        card.line.width = Pt(1)

        ctf = card.text_frame
        ctf.word_wrap = True
        ctf.margin_left = Inches(0.3)
        ctf.margin_right = Inches(0.3)
        ctf.margin_top = Inches(0.2)

        p = ctf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = PRIMARY

        p_sub = ctf.add_paragraph()
        p_sub.text = category
        p_sub.font.size = Pt(10)
        p_sub.font.color.rgb = TEXT_MUTED
        p_sub.space_after = Pt(6)

        for b in bullets:
            pb = ctf.add_paragraph()
            pb.text = f"•  {b}"
            pb.font.size = Pt(10.5)
            pb.font.color.rgb = TEXT_DARK
            pb.space_after = Pt(3)

    # -------------------------------------------------------------
    # SLIDE 4: Frontend Highlights & Live Website
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4, LIGHT_BG)
    add_header(slide4, "4", "Frontend Experience & Live Features", "Clean, responsive, accessible, and themeable user interface")

    # Left: 2 Screenshot previews
    if os.path.exists("docs/screenshots/explore_stays.jpg"):
        slide4.shapes.add_picture("docs/screenshots/explore_stays.jpg", Inches(0.8), Inches(1.8), Inches(5.6), Inches(3.2))
    if os.path.exists("docs/screenshots/my_bookings.jpg"):
        slide4.shapes.add_picture("docs/screenshots/my_bookings.jpg", Inches(0.8), Inches(5.1), Inches(5.6), Inches(1.55))

    # Right: Bullet points card
    card4 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.6), Inches(1.8), Inches(5.9), Inches(4.85))
    card4.fill.solid()
    card4.fill.fore_color.rgb = CARD_LIGHT
    card4.line.color.rgb = BORDER_COLOR
    card4.line.width = Pt(1)

    ctf4 = card4.text_frame
    ctf4.word_wrap = True
    ctf4.margin_left = Inches(0.3)
    ctf4.margin_right = Inches(0.3)
    ctf4.margin_top = Inches(0.3)

    p = ctf4.paragraphs[0]
    p.text = "Key Frontend Highlights"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.space_after = Pt(10)

    fe_bullets = [
        "🔍 Explore & Multi-Filter: Instant debounced search by location or property title with category chips (Mountain, Forest, Riverside, Coastal).",
        "📅 Dynamic Booking Modal: Interactive calendar date picker with date sequence validation & real-time multi-night price calculation.",
        "📊 Host Management Dashboard: Dual-role authorization with dedicated dashboard showing live revenue and occupancy metrics.",
        "🛡️ Production UX & Polish: Global React ErrorBoundary to prevent blank-screen crashes, dark mode support with localStorage persistence."
    ]
    for b in fe_bullets:
        pb = ctf4.add_paragraph()
        pb.text = f"•  {b}"
        pb.font.size = Pt(11)
        pb.font.color.rgb = TEXT_DARK
        pb.space_after = Pt(8)

    # Bottom link bar
    link_bar = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.45))
    link_bar.fill.solid()
    link_bar.fill.fore_color.rgb = PRIMARY
    l_tf = link_bar.text_frame
    lp = l_tf.paragraphs[0]
    lp.text = "🌐 Live Frontend URL: https://trishul-stay-ease.vercel.app"
    lp.font.size = Pt(12)
    lp.font.bold = True
    lp.font.color.rgb = WHITE
    lp.alignment = PP_ALIGN.CENTER

    # -------------------------------------------------------------
    # SLIDE 5: Backend Top 2 APIs & Architecture
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5, LIGHT_BG)
    add_header(slide5, "5", "Backend: Top 2 APIs & Implementation", "Robust asynchronous endpoints with strict input validation and error handling")

    apis = [
        ("API 1: POST /api/ai/travel-plan (AI Eco Planner)", [
            "Purpose: Generates customized sustainable travel itineraries.",
            "Working Mechanism:",
            "  1. Accepts destination, budget, trip duration, guest count, and travel style.",
            "  2. Evaluates in-memory TTL cache to return instant cached responses for duplicate prompts.",
            "  3. Invokes Google Gemini 1.5 Flash with structured JSON system prompts.",
            "  4. Validates schema integrity and gracefully handles API quota boundaries.",
            "Status Code: 200 OK | Rate Limited | Cached"
        ]),
        ("API 2: POST /auth/register & /auth/login (JWT & OAuth Auth)", [
            "Purpose: Provides enterprise-grade authentication & dual role management.",
            "Working Mechanism:",
            "  1. Password hashing with Bcrypt (12 salt rounds) via Passlib.",
            "  2. SlowAPI rate limiting (5 req / 15 min) prevents brute-force credential attacks.",
            "  3. Issues stateless 24-hour HS256 JWT access tokens.",
            "  4. Supports one-click Google OAuth 2.0 token verification.",
            "Status Code: 201 Created / 200 OK | Strict CORS & Security Headers"
        ])
    ]

    for idx, (title, bullets) in enumerate(apis):
        card = slide5.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, 
            Inches(0.8), 
            Inches(1.8 + idx*2.45), 
            Inches(11.7), 
            Inches(2.3)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_LIGHT
        card.line.color.rgb = BORDER_COLOR
        card.line.width = Pt(1)

        ctf = card.text_frame
        ctf.word_wrap = True
        ctf.margin_left = Inches(0.3)
        ctf.margin_top = Inches(0.2)

        p = ctf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13.5)
        p.font.bold = True
        p.font.color.rgb = PRIMARY
        p.space_after = Pt(6)

        for b in bullets:
            pb = ctf.add_paragraph()
            pb.text = b
            pb.font.size = Pt(11)
            pb.font.color.rgb = TEXT_DARK
            pb.space_after = Pt(2)

    # Bottom link bar
    link_bar = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.45))
    link_bar.fill.solid()
    link_bar.fill.fore_color.rgb = CARD_BG
    l_tf = link_bar.text_frame
    lp = l_tf.paragraphs[0]
    lp.text = "📖 Live Swagger UI & Interactive Documentation: https://trishul-stayease.onrender.com/docs"
    lp.font.size = Pt(12)
    lp.font.bold = True
    lp.font.color.rgb = WHITE
    lp.alignment = PP_ALIGN.CENTER

    # -------------------------------------------------------------
    # SLIDE 6: Database & Schema Design
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6, LIGHT_BG)
    add_header(slide6, "6", "Database Selection & Schema Design", "MongoDB Atlas: Fully managed cloud NoSQL database with async Motor driver")

    collections = [
        ("📁 properties Collection", [
            "id: int (Auto-increment PK)",
            "title: string",
            "location: string",
            "price: int (Nightly rate)",
            "type: string (mountain, forest...)",
            "status: available | booked",
            "amenities: array of strings",
            "images: array of URLs",
            "host_id: string"
        ]),
        ("👤 users Collection", [
            "id: ObjectId (PK)",
            "full_name: string",
            "email: string (Unique Index)",
            "phone: string (Unique)",
            "password_hash: string",
            "role: guest | host",
            "avatar_url: string",
            "created_at: timestamp"
        ]),
        ("📅 bookings Collection", [
            "id: int (PK sequence)",
            "property_id: int",
            "guest_id: string",
            "host_id: string",
            "start_date: date",
            "end_date: date",
            "total_price: int",
            "status: upcoming | completed | cancelled"
        ]),
        ("🔢 counters & wishlists", [
            "counters: Atomic auto-increment sequence generation via $inc operator.",
            "wishlists: Stores user favorite property ID arrays for zero-friction access.",
            "Indexes: Compound indexes on (status, price) and unique text search on (title, location)."
        ])
    ]

    for idx, (title, bullets) in enumerate(collections):
        col = idx % 2
        row = idx // 2
        card = slide6.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, 
            Inches(0.8 + col*5.95), 
            Inches(1.8 + row*2.45), 
            Inches(5.75), 
            Inches(2.25)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_LIGHT
        card.line.color.rgb = BORDER_COLOR
        card.line.width = Pt(1)

        ctf = card.text_frame
        ctf.word_wrap = True
        ctf.margin_left = Inches(0.3)
        ctf.margin_top = Inches(0.2)

        p = ctf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = PRIMARY
        p.space_after = Pt(4)

        for b in bullets:
            pb = ctf.add_paragraph()
            pb.text = f"•  {b}"
            pb.font.size = Pt(10)
            pb.font.color.rgb = TEXT_DARK
            pb.space_after = Pt(2)

    # Bottom rationale bar
    bar = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.45))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    b_tf = bar.text_frame
    bp = b_tf.paragraphs[0]
    bp.text = "💡 Why MongoDB? Flexible schema matches polymorphic eco-stay amenities + zero-downtime horizontal scaling."
    bp.font.size = Pt(11)
    bp.font.bold = True
    bp.font.color.rgb = WHITE
    bp.alignment = PP_ALIGN.CENTER

    # -------------------------------------------------------------
    # SLIDE 7: AI Feature & LLM Integration
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7, LIGHT_BG)
    add_header(slide7, "7", "AI Feature: Eco Travel Planner", "Harnessing Google Gemini 1.5 Flash for personalized sustainable itinerary generation")

    # Left: AI Planner Screenshot
    if os.path.exists("docs/screenshots/ai_planner.jpg"):
        slide7.shapes.add_picture("docs/screenshots/ai_planner.jpg", Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.85))

    # Right: AI Details Card
    card7 = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.6), Inches(1.8), Inches(5.9), Inches(4.85))
    card7.fill.solid()
    card7.fill.fore_color.rgb = CARD_LIGHT
    card7.line.color.rgb = BORDER_COLOR
    card7.line.width = Pt(1)

    ctf7 = card7.text_frame
    ctf7.word_wrap = True
    ctf7.margin_left = Inches(0.3)
    ctf7.margin_right = Inches(0.3)
    ctf7.margin_top = Inches(0.3)

    p = ctf7.paragraphs[0]
    p.text = "Google Gemini 1.5 Flash Integration"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.space_after = Pt(8)

    ai_bullets = [
        "🤖 Model & SDK: Google Gemini 1.5 Flash via official `@google/genai` Python SDK for fast token latency and strict structured JSON schema.",
        "🎯 Use Case: Instant day-by-day eco itineraries tailored to budget, days, guest count, and travel style with packing & green guidelines.",
        "⚡ Server-Side TTL Cache: In-memory dictionary cache returns instant responses for repeat destination queries, preventing API quota exhaustion.",
        "🛡️ Schema Validation: Pydantic validation handles JSON formatting errors and markdown codeblock stripping gracefully."
    ]
    for b in ai_bullets:
        pb = ctf7.add_paragraph()
        pb.text = f"•  {b}"
        pb.font.size = Pt(11)
        pb.font.color.rgb = TEXT_DARK
        pb.space_after = Pt(8)

    # -------------------------------------------------------------
    # SLIDE 8: Hosting & Cloud Infrastructure
    # -------------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8, LIGHT_BG)
    add_header(slide8, "8 & 9", "Cloud Hosting Services & All Live URLs", "Production-grade multi-cloud architecture deployed on Vercel, Render & MongoDB Atlas")

    hosts = [
        ("▲ Vercel (Frontend Hosting)", [
            "Role: Hosts the React 19 + Vite Single Page Application.",
            "Features: Global CDN Edge Network, automated preview deployments on git push, automatic gzip/brotli compression.",
            "Live URL: https://trishul-stay-ease.vercel.app"
        ]),
        ("⚡ Render (Backend Hosting)", [
            "Role: Hosts the FastAPI Python 3.11 ASGI Web Service.",
            "Features: Managed Linux container, automated build pipelines from GitHub, secure environment secret injection.",
            "Live URL: https://trishul-stayease.onrender.com"
        ]),
        ("🍃 MongoDB Atlas (Cloud Database)", [
            "Role: Fully managed cloud database cluster.",
            "Features: Automatic daily backups, connection pooling, and SSL-encrypted in-transit data access.",
            "Cluster: M0 Free Sandbox (AWS Mumbai Region)"
        ]),
        ("🐙 GitHub (Source Code Repository)", [
            "Role: Complete open-source version controlled repository.",
            "Features: Modular folder structure, Postman collection, atomic commits, comprehensive README.",
            "Repository: https://github.com/Riddhi1204/Trishul_StayEase"
        ])
    ]

    for idx, (title, bullets) in enumerate(hosts):
        col = idx % 2
        row = idx // 2
        card = slide8.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, 
            Inches(0.8 + col*5.95), 
            Inches(1.8 + row*2.45), 
            Inches(5.75), 
            Inches(2.25)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_LIGHT
        card.line.color.rgb = BORDER_COLOR
        card.line.width = Pt(1)

        ctf = card.text_frame
        ctf.word_wrap = True
        ctf.margin_left = Inches(0.3)
        ctf.margin_top = Inches(0.2)

        p = ctf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = PRIMARY
        p.space_after = Pt(4)

        for b in bullets:
            pb = ctf.add_paragraph()
            pb.text = f"•  {b}"
            pb.font.size = Pt(10.5)
            pb.font.color.rgb = TEXT_DARK
            pb.space_after = Pt(3)

    # Bottom bar summary
    bar = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.45))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    b_tf = bar.text_frame
    bp = b_tf.paragraphs[0]
    bp.text = "🔗 All links are live, SSL secured (HTTPS), and publicly accessible without restrictions."
    bp.font.size = Pt(11)
    bp.font.bold = True
    bp.font.color.rgb = WHITE
    bp.alignment = PP_ALIGN.CENTER

    # -------------------------------------------------------------
    # SLIDE 9: Internship Reflection & Key Learnings
    # -------------------------------------------------------------
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9, LIGHT_BG)
    add_header(slide9, "10", "Reflection on the Internship & Learnings", "A transformative journey building a complete full-stack web product from scratch")

    reflections = [
        ("🚀 Technical Growth & Mastery", [
            "Full-Stack Integration: Built end-to-end synergy between React SPA, async FastAPI, and MongoDB Atlas.",
            "Authentication & Security: Implemented production JWT authentication, Bcrypt password hashing, and OAuth 2.0.",
            "AI Engineering: Gained practical expertise integrating Google Gemini LLM with structured prompts and caching.",
            "Performance Optimization: Leveraged React `useMemo`, debounce searches, and connection pooling."
        ]),
        ("💼 Engineering Best Practices & Experience", [
            "Defensive Coding: Handled edge cases with React Error Boundaries, Pydantic validations, and CORS middlewares.",
            "Agile Version Control: Maintained clean, atomic git commit histories following conventional commits standard.",
            "Production Cloud Deployment: Managed live multi-cloud environments across Vercel, Render, and Atlas.",
            "Overall Experience: Incredible hands-on exposure solving real-world eco-tourism booking challenges."
        ])
    ]

    for idx, (title, bullets) in enumerate(reflections):
        card = slide9.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, 
            Inches(0.8 + idx*5.95), 
            Inches(1.8), 
            Inches(5.75), 
            Inches(4.8)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_LIGHT
        card.line.color.rgb = BORDER_COLOR
        card.line.width = Pt(1)

        ctf = card.text_frame
        ctf.word_wrap = True
        ctf.margin_left = Inches(0.3)
        ctf.margin_right = Inches(0.3)
        ctf.margin_top = Inches(0.3)

        p = ctf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = PRIMARY
        p.space_after = Pt(12)

        for b in bullets:
            pb = ctf.add_paragraph()
            pb.text = f"•  {b}"
            pb.font.size = Pt(11.5)
            pb.font.color.rgb = TEXT_DARK
            pb.space_after = Pt(8)

    # -------------------------------------------------------------
    # SLIDE 10: Conclusion / Thank You (Dark Theme)
    # -------------------------------------------------------------
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide10, DARK_BG)

    tb = slide10.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10.3), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Thank You! 🙏"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Trishul StayEase — Sustainable Homestays, Direct Bookings"
    p2.font.size = Pt(18)
    p2.font.color.rgb = ACCENT
    p2.alignment = PP_ALIGN.CENTER
    p2.space_after = Pt(20)

    p3 = tf.add_paragraph()
    p3.text = "Candidate: Riddhi Kumari  |  Intern ID: 26100462"
    p3.font.size = Pt(14)
    p3.font.color.rgb = WHITE
    p3.alignment = PP_ALIGN.CENTER

    p4 = tf.add_paragraph()
    p4.text = "🌐 Live App: https://trishul-stay-ease.vercel.app"
    p4.font.size = Pt(14)
    p4.font.color.rgb = GOLD
    p4.alignment = PP_ALIGN.CENTER

    p5 = tf.add_paragraph()
    p5.text = "🐙 GitHub: https://github.com/Riddhi1204/Trishul_StayEase"
    p5.font.size = Pt(14)
    p5.font.color.rgb = RGBColor(180, 205, 195)
    p5.alignment = PP_ALIGN.CENTER

    output_path = "Trishul_StayEase_Final_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

if __name__ == "__main__":
    create_presentation()
